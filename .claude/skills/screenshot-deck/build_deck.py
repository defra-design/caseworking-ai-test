# Build the documentation deck for Grasslands caselist + user-management changes.
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

SHOTS = r"C:\Users\chasl\Documents\ai-driven-main\_ppt_shots"
OUT = r"C:\Users\chasl\Downloads\grasslands-caselist-and-user-management.pptx"

GREEN = RGBColor(0x00, 0x70, 0x3C)
INK = RGBColor(0x0B, 0x0C, 0x0C)
MUTED = RGBColor(0x50, 0x5A, 0x5F)
LINE = RGBColor(0xB1, 0xB4, 0xB6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PANEL = RGBColor(0xF3, 0xF2, 0xF1)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]


def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=4):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, (text, size, color, bold, bullet) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        p.space_before = Pt(0)
        if bullet:
            p.line_spacing = 1.05
        r = p.add_run()
        r.text = (("•  " + text) if bullet else text)
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = "Arial"
    return tb


def kicker(slide, x, y, text):
    # small green square motif + section label
    sq = slide.shapes.add_shape(1, Inches(x), Inches(y + 0.03), Inches(0.16), Inches(0.16))
    sq.fill.solid(); sq.fill.fore_color.rgb = GREEN; sq.line.fill.background()
    add_text(slide, x + 0.28, y, 9, 0.35,
             [(text.upper(), 13, GREEN, True, False)])


def fit(img_path, bx, by, bw, bh, valign="top", halign="center"):
    iw, ih = Image.open(img_path).size
    ar = iw / ih
    if (bw / bh) > ar:
        h = bh; w = bh * ar
    else:
        w = bw; h = bw / ar
    x = bx + {"left": 0, "center": (bw - w) / 2, "right": bw - w}[halign]
    y = by + {"top": 0, "center": (bh - h) / 2, "bottom": bh - h}[valign]
    return x, y, w, h


def add_image(slide, img_path, bx, by, bw, bh, valign="top", halign="center"):
    x, y, w, h = fit(img_path, bx, by, bw, bh, valign, halign)
    pic = slide.shapes.add_picture(img_path, Inches(x), Inches(y), Inches(w), Inches(h))
    pic.line.color.rgb = LINE
    pic.line.width = Pt(0.75)
    return pic


def content_slide(section, title, why, bullets, images, image_side="right", captions=None):
    s = prs.slides.add_slide(BLANK)
    M = 0.6
    kicker(s, M, 0.5, section)
    add_text(s, M, 0.92, SW - 2 * M, 1.0, [(title, 26, INK, True, False)])

    top = 2.0
    bottom = 7.0
    if len(images) == 1:
        img_w = 6.4
        img_bottom = 6.85  # clear bottom margin so tall screenshots don't touch the slide edge
        if image_side == "right":
            img_box = (SW - M - img_w, top, img_w, img_bottom - top)
            txt_x = M; txt_w = SW - 2 * M - img_w - 0.5
        else:
            img_box = (M, top, img_w, img_bottom - top)
            txt_x = M + img_w + 0.5; txt_w = SW - 2 * M - img_w - 0.5
        add_image(s, images[0], *img_box, valign="top")
        runs = [(why, 14, MUTED, False, False), ("", 6, MUTED, False, False)]
        runs += [(b, 14, INK, False, True) for b in bullets]
        add_text(s, txt_x, top + 0.05, txt_w, bottom - top, runs, space=7)
    else:
        # two images across the top, bullets below
        gap = 0.6
        iw = (SW - 2 * M - gap) / 2
        ih = 2.45
        for idx, img in enumerate(images[:2]):
            bx = M + idx * (iw + gap)
            add_image(s, img, bx, top, iw, ih, valign="top", halign="center")
            if captions:
                add_text(s, bx, top + ih + 0.06, iw, 0.3,
                         [(captions[idx], 11, MUTED, True, False)], align=PP_ALIGN.CENTER)
        ty = top + ih + 0.5
        runs = [(why, 13.5, MUTED, False, False)]
        runs += [(b, 13, INK, False, True) for b in bullets]
        add_text(s, M, ty, SW - 2 * M, bottom - ty, runs, space=4)
    return s


def p(path):
    return os.path.join(SHOTS, path)


# ---- Title slide ----
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = GREEN; bg.line.fill.background()
bg.shadow.inherit = False
s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
add_text(s, 0.9, 2.4, 11.5, 2.0, [
    ("Grasslands caselist & user management", 40, WHITE, True, False),
    ("Documenting the changes — features and the reasons behind them", 20, WHITE, False, False),
], space=10)
add_text(s, 0.9, 6.4, 11.5, 0.5, [
    ("DEFRA — Manage rural grant applications  ·  19 June 2026", 13, WHITE, False, False)])

# ---- Content slides ----
content_slide(
    "Search and filtering",
    "Find a case fast in a growing queue",
    "With 85+ cases spread across teams, caseworkers need to locate one case and narrow a long list.",
    [
        "Collapsible “Find an application” panel keeps the list clean",
        "Search by application ID, business, SBI, status or assignee",
        "Filter by Assignee and by Status (checkbox groups)",
        "Live counts beside every option, reconciled to the data",
        "Statuses limited to the agreed FRPS-D2 set",
    ],
    [p("filters_open.png")], image_side="right",
)

content_slide(
    "Multiple assignment",
    "Assign or reassign in bulk, with an audit trail",
    "Managers move work between caseworkers and clear unallocated cases without editing one case at a time.",
    [
        "Select checkbox per row + “Reassign selected cases” (bulk)",
        "One pattern: team → caseworker → reason → confirm",
        "Caseworker picker filtered to the team, showing current load",
        "Specialist tasks filter the picker to specialism holders",
        "Every change logs assignment history and a Timeline event",
    ],
    [p("um_allocate.png")], image_side="left",
)

content_slide(
    "Tab structure and teams",
    "My cases, team context and completed — in three tabs",
    "A caseworker's own work, their team's queue, and finished cases each need a clear home.",
    [
        "Tabs: My cases  ·  [team / All]  ·  Completed cases",
        "“Switch team” selector (Team A/B/C + All cases) above the tabs",
        "Middle tab shows the selected team, or everyone",
        "Completed split out (Agreement accepted / Rejected / Withdrawn)",
        "My and team tabs show active work only",
    ],
    [p("tabs_teams.png")], image_side="right",
)

content_slide(
    "Tab structure and teams",
    "Teams are now real, managed entities",
    "The caselist's teams were hardcoded in a data file — user management now owns them.",
    [
        "Teams list: lead, members, open cases, unallocated",
        "Three work-allocation teams, three caseworkers each",
        "Per-team workload visible at a glance",
        "Unallocated work surfaced for triage",
        "Built to the agreed team-structure design brief",
    ],
    [p("um_teams.png")], image_side="left",
)

content_slide(
    "Context persistence across tabs",
    "Pick a context once — it follows you across tabs",
    "Switching team and then changing tab should not reset the view.",
    [
        "Context (a team, or All cases) is held in the session",
        "Switch to All cases → the context tab shows everyone…",
        "…and the Completed tab shows all completed too",
        "Pagination preserves context; an invalid one falls back safely",
    ],
    [p("context_all.png"), p("completed_all.png")],
    captions=["Context = All cases", "Completed tab keeps All cases"],
)

content_slide(
    "User management",
    "User management, iterated for teams (UM2)",
    "The team caselist needs users that belong to teams — not just scheme permissions.",
    [
        "Landing now offers Users · Teams · Scheme roles",
        "“Roles” renamed “Scheme roles” to remove ambiguity",
        "Users list gains Teams, Access role, Specialisms, Status",
        "Access role = Caseworker or Manager; multi-team membership",
    ],
    [p("um_landing.png"), p("um_users.png")],
    captions=["New landing — three areas", "Users list — new columns"],
)

content_slide(
    "User management",
    "A user is a team member with a role and specialisms",
    "This is what turns a user into a caselist caseworker.",
    [
        "Teams (multi-select) — a user can be in several",
        "Access role (Caseworker by default / Manager + scope)",
        "Specialisms (Finance officer, Quality control)",
        "Scheme roles kept separate and unchanged",
    ],
    [p("um_create.png")], image_side="right",
)

content_slide(
    "User management",
    "Team detail: membership and workload",
    "Managers need to balance load and never orphan work.",
    [
        "Members shown with open-case count and oldest-case age",
        "Set the team lead; add or remove members",
        "Removing someone with open cases starts a reassign step",
        "Can't delete a team that still owns open cases",
    ],
    [p("um_team.png")], image_side="left",
)

# ---- Summary ----
s = prs.slides.add_slide(BLANK)
M = 0.6
kicker(s, M, 0.5, "Summary")
add_text(s, M, 0.92, SW - 2 * M, 1.0, [("What changed, and why", 26, INK, True, False)])
add_text(s, M, 2.0, SW - 2 * M, 4.6, [
    ("Search and filtering — find and narrow a large queue.", 15, INK, False, True),
    ("Multiple assignment — bulk reassign via one audited team → caseworker pattern.", 15, INK, False, True),
    ("Tab structure and teams — My / team-or-All / Completed, with real, managed teams.", 15, INK, False, True),
    ("Context persistence — the chosen team/All view follows the user across tabs.", 15, INK, False, True),
    ("User management — teams, access roles and specialisms added to make the caselist run on real users.", 15, INK, False, True),
    ("", 8, MUTED, False, False),
    ("These are working prototype iterations. Next: build to the resolved design brief (MVP first).", 13, MUTED, False, False),
], space=10)

prs.save(OUT)
print("saved", OUT, "slides:", len(prs.slides._sldIdLst))
